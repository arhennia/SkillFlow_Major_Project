"""
generate_presentation.py
=========================
Builds a crisp 7-slide PowerPoint deck
(`presentation/Churn_Prediction_Presentation.pptx`) summarizing the
Customer Churn Prediction & Intelligent Retention System project,
using python-pptx.

Reads from models/:
    - model_comparison.csv
    - champion_model_info.json

Embeds plots produced by build_notebook.py (saved in notebooks/):
    - plot_01_churn_distribution.png
    - plot_04_churn_by_contract_payment.png
    - plot_09_roc_curves.png
    - plot_11_feature_importance.png

Run (from project root):
    python src/generate_presentation.py
"""

import json
import os

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

_SRC_DIR = os.path.dirname(os.path.abspath(__file__))
_PROJECT_ROOT = os.path.dirname(_SRC_DIR)

MODELS_DIR = os.path.join(_PROJECT_ROOT, "models")
PLOTS_DIR = os.path.join(_PROJECT_ROOT, "notebooks")
PRESENTATION_DIR = os.path.join(_PROJECT_ROOT, "presentation")
os.makedirs(PRESENTATION_DIR, exist_ok=True)

OUTPUT_PATH = os.path.join(PRESENTATION_DIR, "Churn_Prediction_Presentation.pptx")

# ---------------------------------------------------------------------
# Load supporting data
# ---------------------------------------------------------------------
comparison_df = pd.read_csv(os.path.join(MODELS_DIR, "model_comparison.csv"))
with open(os.path.join(MODELS_DIR, "champion_model_info.json")) as f:
    champion_info = json.load(f)
champion_name = champion_info["champion_model"]
metrics = champion_info["metrics"]

# ---------------------------------------------------------------------
# Palette & constants
# ---------------------------------------------------------------------
NAVY = RGBColor(0x1F, 0x38, 0x64)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
ORANGE = RGBColor(0xDD, 0x84, 0x52)
DARK_GREY = RGBColor(0x40, 0x40, 0x40)
MID_GREY = RGBColor(0x70, 0x70, 0x70)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_background(slide, color=WHITE):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=18, color=DARK_GREY,
                 bold=False, italic=False, align=PP_ALIGN.LEFT, font="Calibri",
                 line_spacing=1.15, anchor=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return box


def add_bullets(slide, left, top, width, height, items, size=15, color=DARK_GREY,
                 bold_lead=False, space_after=10, font="Calibri", bullet_color=ACCENT_BLUE):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        p.line_spacing = 1.12
        # bullet marker as separate run for color control
        marker = p.add_run()
        marker.text = "●  "
        marker.font.size = Pt(size - 3)
        marker.font.color.rgb = bullet_color
        marker.font.name = font
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run()
            r1.text = lead
            r1.font.bold = True
            r1.font.size = Pt(size)
            r1.font.color.rgb = color
            r1.font.name = font
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
            r2.font.name = font
        else:
            r = p.add_run()
            r.text = item
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.name = font
    return box


def add_title_bar(slide, title, subtitle=None, dark=False):
    """Slide header: bold title + optional grey subtitle."""
    title_color = WHITE if dark else NAVY
    add_textbox(slide, Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.7),
                title, size=28, color=title_color, bold=True)
    if subtitle:
        sub_color = RGBColor(0xC7, 0xD5, 0xE8) if dark else MID_GREY
        add_textbox(slide, Inches(0.6), Inches(0.95), Inches(12.1), Inches(0.4),
                    subtitle, size=14, color=sub_color, italic=True)


def add_footer(slide, page_num, dark=False):
    color = RGBColor(0xB0, 0xB8, 0xC4) if dark else RGBColor(0xA0, 0xA0, 0xA0)
    add_textbox(slide, Inches(0.6), Inches(7.12), Inches(8), Inches(0.3),
                "Customer Churn Prediction & Intelligent Retention System",
                size=9, color=color)
    add_textbox(slide, Inches(12.3), Inches(7.12), Inches(0.5), Inches(0.3),
                str(page_num), size=9, color=color, align=PP_ALIGN.RIGHT)


def stat_card(slide, left, top, width, height, value, label, value_color=NAVY):
    box = slide.shapes.add_shape(1, left, top, width, height)  # 1 = MSO_SHAPE.RECTANGLE
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = RGBColor(0xE0, 0xE4, 0xEA)
    box.line.width = Pt(0.75)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = value
    r1.font.size = Pt(30)
    r1.font.bold = True
    r1.font.color.rgb = value_color
    r1.font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(12.5)
    r2.font.color.rgb = MID_GREY
    r2.font.name = "Calibri"
    return box


def add_picture_fit(slide, path, left, top, max_w, max_h):
    """Add an image, scaled to fit within (max_w, max_h) preserving aspect ratio, centered."""
    from PIL import Image
    with Image.open(path) as im:
        iw, ih = im.size
    ratio = min(max_w / iw, max_h / ih)
    w, h = int(iw * ratio), int(ih * ratio)
    left_c = left + (max_w - w) // 2
    top_c = top + (max_h - h) // 2
    slide.shapes.add_picture(path, left_c, top_c, width=Emu(w), height=Emu(h))


# =====================================================================
# SLIDE 1 — Title & Executive Overview
# =====================================================================
s = add_slide()
set_background(s, NAVY)
add_textbox(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.3),
            "Customer Churn Prediction &", size=40, color=WHITE, bold=True)
add_textbox(s, Inches(0.9), Inches(3.05), Inches(11.5), Inches(1.3),
            "Intelligent Retention System", size=40, color=WHITE, bold=True)
add_textbox(s, Inches(0.9), Inches(4.0), Inches(11), Inches(0.6),
            "An End-to-End Machine Learning Pipeline for Proactive Customer Retention",
            size=18, color=RGBColor(0xC7, 0xD5, 0xE8), italic=True)
add_textbox(s, Inches(0.9), Inches(6.5), Inches(11), Inches(0.5),
            "Data Science & Machine Learning Internship Project  |  IBM Telco Customer Churn Dataset",
            size=13, color=RGBColor(0x9F, 0xB0, 0xCC))

# =====================================================================
# SLIDE 2 — Business Problem & Financial Impact of Churn
# =====================================================================
s = add_slide()
set_background(s)
add_title_bar(s, "Business Problem & Financial Impact",
              "Why churn prediction is a high-leverage lever for the business")

stat_card(s, Inches(0.6), Inches(1.55), Inches(2.7), Inches(1.5), "26.5%", "Overall Churn Rate")
stat_card(s, Inches(3.45), Inches(1.55), Inches(2.7), Inches(1.5), "7,043", "Customers Analyzed")
stat_card(s, Inches(6.3), Inches(1.55), Inches(2.7), Inches(1.5), "1,869", "Customers Lost to Churn")
stat_card(s, Inches(9.15), Inches(1.55), Inches(3.15), Inches(1.5), "5-25x", "Cost to Acquire vs. Retain*")

add_bullets(s, Inches(0.6), Inches(3.35), Inches(11.8), Inches(3.3), [
    ("The revenue problem: ", "roughly 1 in 4 customers churns, representing a continuous, "
     "compounding drain on recurring subscription revenue."),
    ("The acquisition-cost problem: ", "replacing a lost customer costs significantly more than "
     "retaining one — every churned customer is a doubly expensive loss."),
    ("The status-quo problem: ", "without predictive scoring, retention efforts are either "
     "wasted on blanket-wide offers or applied too late, after a customer has already decided to leave."),
    ("The opportunity: ", "a model that reliably flags at-risk customers early lets the business "
     "target retention spend where it has the highest return."),
], size=16, space_after=16)

add_textbox(s, Inches(0.6), Inches(6.85), Inches(11.5), Inches(0.3),
            "*Widely cited industry benchmark for customer acquisition vs. retention cost.",
            size=10, color=MID_GREY, italic=True)
add_footer(s, 2)

# =====================================================================
# SLIDE 3 — EDA Key Insights
# =====================================================================
s = add_slide()
set_background(s)
add_title_bar(s, "Exploratory Data Analysis — Key Insights",
              "What the data reveals about who churns, and why")

plot1 = os.path.join(PLOTS_DIR, "plot_01_churn_distribution.png")
plot4 = os.path.join(PLOTS_DIR, "plot_04_churn_by_contract_payment.png")
if os.path.exists(plot1):
    add_picture_fit(s, plot1, Inches(0.5), Inches(1.55), Inches(3.6), Inches(4.9))
if os.path.exists(plot4):
    add_picture_fit(s, plot4, Inches(4.3), Inches(1.55), Inches(8.5), Inches(4.9))

add_bullets(s, Inches(0.6), Inches(6.55), Inches(11.8), Inches(0.7), [
    ("Contract type dominates: ", "month-to-month customers churn far more than one- or two-year "
     "contract holders — the single strongest churn signal in the data."),
], size=13.5, space_after=0)
add_footer(s, 3)

# =====================================================================
# SLIDE 4 — ML Pipeline & Feature Engineering
# =====================================================================
s = add_slide()
set_background(s)
add_title_bar(s, "Machine Learning Pipeline & Feature Engineering",
              "From raw data to a production-ready feature matrix")

stages = [
    ("1. Data Hygiene", "Coerce TotalCharges to numeric; impute 11 blank-string values "
     "(new customers, tenure = 0) as 0."),
    ("2. Feature Engineering", "Create Total_Services_Used, Estimated_LTV, Tenure_Group, "
     "and Payment_Risk_Score."),
    ("3. Encoding & Scaling", "One-hot encode nominal features; standard-scale continuous "
     "features (tenure, charges, LTV)."),
    ("4. Stratified Split", "80/20 train-test split, stratified on Churn, random_state=42, "
     "to preserve class balance."),
]

card_w = Inches(2.95)
gap = Inches(0.15)
left0 = Inches(0.6)
top0 = Inches(1.75)
for i, (title, desc) in enumerate(stages):
    left = left0 + i * (card_w + gap)
    box = s.shapes.add_shape(1, left, top0, card_w, Inches(2.6))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = RGBColor(0xE0, 0xE4, 0xEA)
    box.line.width = Pt(0.75)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(14)
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.bold = True
    r1.font.size = Pt(15)
    r1.font.color.rgb = ACCENT_BLUE
    r1.font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.space_before = Pt(8)
    p2.line_spacing = 1.15
    r2 = p2.add_run()
    r2.text = desc
    r2.font.size = Pt(12.5)
    r2.font.color.rgb = DARK_GREY
    r2.font.name = "Calibri"

add_textbox(s, Inches(0.6), Inches(4.65), Inches(11.8), Inches(0.4),
            "Engineered Features", size=17, color=NAVY, bold=True)
add_bullets(s, Inches(0.6), Inches(5.1), Inches(11.8), Inches(1.9), [
    ("Total_Services_Used — ", "count of subscribed add-on services (0-6); higher counts signal engagement."),
    ("Estimated_LTV — ", "MonthlyCharges × tenure; a proxy for cumulative customer value."),
    ("Tenure_Group — ", "lifecycle-stage bucket (0-12 / 12-24 / 24-48 / 48+ months)."),
    ("Payment_Risk_Score — ", "flags non-automatic payment methods as higher churn risk."),
], size=13.5, space_after=6)
add_footer(s, 4)

# =====================================================================
# SLIDE 5 — Model Comparison & Evaluation Metrics
# =====================================================================
s = add_slide()
set_background(s)
add_title_bar(s, "Model Comparison & Evaluation Metrics",
              f"Champion model: {champion_name}  |  ROC-AUC {metrics['ROC-AUC']:.3f}  |  Recall {metrics['Recall']:.1%}")

# Table
rows, cols = comparison_df.shape[0] + 1, comparison_df.shape[1]
tbl_left, tbl_top, tbl_w, tbl_h = Inches(0.6), Inches(1.7), Inches(6.6), Inches(3.5)
table_shape = s.shapes.add_table(rows, cols, tbl_left, tbl_top, tbl_w, tbl_h)
table = table_shape.table

for j, col_name in enumerate(comparison_df.columns):
    cell = table.cell(0, j)
    cell.text = col_name
    cell.fill.solid()
    cell.fill.fore_color.rgb = NAVY
    for p in cell.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.bold = True
            r.font.size = Pt(11)
            r.font.color.rgb = WHITE

champ_idx = comparison_df.index[comparison_df["Model"] == champion_name][0]
for i, row in comparison_df.iterrows():
    for j, val in enumerate(row):
        cell = table.cell(i + 1, j)
        cell.text = f"{val:.4f}" if isinstance(val, float) else str(val)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0xD9, 0xE8, 0xF5) if i == champ_idx else WHITE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            for r in p.runs:
                r.font.size = Pt(10.5)
                r.font.bold = (i == champ_idx)
                r.font.color.rgb = DARK_GREY

plot9 = os.path.join(PLOTS_DIR, "plot_09_roc_curves.png")
if os.path.exists(plot9):
    add_picture_fit(s, plot9, Inches(7.4), Inches(1.6), Inches(5.4), Inches(3.7))

add_bullets(s, Inches(0.6), Inches(5.5), Inches(11.8), Inches(1.5), [
    ("Priority metrics: ", "Recall and ROC-AUC are prioritized over raw Accuracy — a missed "
     "churner (False Negative) costs far more than a wasted retention offer (False Positive)."),
    ("Champion selection: ", f"{champion_name} achieves the best balance of ROC-AUC and Recall "
     "after 5-fold Stratified cross-validated hyperparameter tuning."),
], size=13.5, space_after=8)
add_footer(s, 5)

# =====================================================================
# SLIDE 6 — Strategic Retention Playbook
# =====================================================================
s = add_slide()
set_background(s)
add_title_bar(s, "Strategic Retention Playbook",
              "Five actionable strategies mapped directly to model findings")

strategies = [
    ("Personalized Offers", "Target high-risk customers (70%+ churn probability) with tailored "
     "discounts, credits, or upgrades instead of generic promotions."),
    ("Loyalty Rewards", "Incentivize migration from month-to-month to annual contracts — the "
     "single strongest churn driver identified by the model."),
    ("Early Churn Alerts", "Score every customer monthly; automatically alert the retention "
     "team when risk crosses into the medium/high band."),
    ("Engagement Campaigns", "Invest in structured onboarding during the first 12 months, "
     "when churn risk is highest."),
    ("Upgrade Incentives", "Bundle add-on services (security, tech support) at a discount to "
     "raise switching costs and perceived value."),
]

card_w = Inches(2.28)
gap = Inches(0.13)
left0 = Inches(0.6)
top0 = Inches(1.8)
for i, (title, desc) in enumerate(strategies):
    left = left0 + i * (card_w + gap)
    box = s.shapes.add_shape(1, left, top0, card_w, Inches(4.6))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = RGBColor(0xE0, 0xE4, 0xEA)
    box.line.width = Pt(0.75)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(16)

    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run()
    r0.text = str(i + 1)
    r0.font.size = Pt(26)
    r0.font.bold = True
    r0.font.color.rgb = ACCENT_BLUE
    r0.font.name = "Calibri"

    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.CENTER
    p1.space_before = Pt(10)
    r1 = p1.add_run()
    r1.text = title
    r1.font.bold = True
    r1.font.size = Pt(14)
    r1.font.color.rgb = NAVY
    r1.font.name = "Calibri"

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)
    p2.line_spacing = 1.15
    r2 = p2.add_run()
    r2.text = desc
    r2.font.size = Pt(11)
    r2.font.color.rgb = DARK_GREY
    r2.font.name = "Calibri"

add_footer(s, 6)

# =====================================================================
# SLIDE 7 — Conclusion & Implementation Roadmap
# =====================================================================
s = add_slide()
set_background(s, NAVY)
add_textbox(s, Inches(0.6), Inches(0.5), Inches(11.5), Inches(0.7),
            "Conclusion & Implementation Roadmap", size=28, color=WHITE, bold=True)

add_bullets(s, Inches(0.6), Inches(1.6), Inches(11.8), Inches(2.4), [
    (f"{champion_name} is production-ready: ", f"ROC-AUC of {metrics['ROC-AUC']:.3f} and Recall of "
     f"{metrics['Recall']:.1%} on held-out data, correctly identifying roughly 4 in 5 at-risk customers."),
    ("Findings map directly to action: ", "contract type, tenure, and payment method are both the "
     "top predictive features and the top retention levers."),
    ("The system closes the loop: ", "predict risk → intervene proactively → measure impact on "
     "churn rate over time."),
], size=15.5, color=RGBColor(0xE4, 0xEA, 0xF5), bullet_color=ORANGE, space_after=14)

add_textbox(s, Inches(0.6), Inches(4.15), Inches(6), Inches(0.4),
            "Next Steps", size=18, color=WHITE, bold=True)

roadmap = [
    "Deploy the scoring pipeline into the CRM for real-time customer risk flags",
    "A/B test retention offers against model-identified risk segments",
    "Integrate Early Churn Alerts into the retention team's workflow",
    "Retrain the model periodically as customer behavior evolves",
]
add_bullets(s, Inches(0.6), Inches(4.65), Inches(11.8), Inches(2.2), roadmap,
            size=14.5, color=RGBColor(0xE4, 0xEA, 0xF5), bullet_color=ORANGE, space_after=10)

add_textbox(s, Inches(0.6), Inches(7.0), Inches(11.8), Inches(0.4),
            "Thank you  |  Questions & Discussion", size=14,
            color=RGBColor(0x9F, 0xB0, 0xCC), italic=True, align=PP_ALIGN.CENTER)

prs.save(OUTPUT_PATH)
print(f"Presentation generated: {OUTPUT_PATH}")
